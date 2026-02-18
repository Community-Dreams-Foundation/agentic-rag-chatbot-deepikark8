"""
Document Processor - Handles multiple file formats
"""
import os
from typing import List, Dict
from pypdf import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation


class DocumentProcessor:
    
    @staticmethod
    def process_pdf(filepath: str) -> List[Dict]:
        """Extract text from PDF"""
        documents = []
        reader = PdfReader(filepath)
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text.strip():
                documents.append({
                    'text': text,
                    'metadata': {
                        'source': os.path.basename(filepath),
                        'page': i,
                        'type': 'pdf'
                    }
                })
        return documents
    
    @staticmethod
    def process_docx(filepath: str) -> List[Dict]:
        """Extract text from Word document"""
        doc = Document(filepath)
        text = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
        
        return [{
            'text': text,
            'metadata': {
                'source': os.path.basename(filepath),
                'type': 'docx'
            }
        }]
    
    @staticmethod
    def process_txt(filepath: str) -> List[Dict]:
        """Extract text from TXT file"""
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
        
        return [{
            'text': text,
            'metadata': {
                'source': os.path.basename(filepath),
                'type': 'txt'
            }
        }]
    
    @staticmethod
    def process_xlsx(filepath: str) -> List[Dict]:
        """Extract text from Excel file"""
        wb = openpyxl.load_workbook(filepath, data_only=True)
        documents = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    rows.append(row_text)
            
            text = '\n'.join(rows)
            if text.strip():
                documents.append({
                    'text': text,
                    'metadata': {
                        'source': os.path.basename(filepath),
                        'sheet': sheet_name,
                        'type': 'xlsx'
                    }
                })
        
        return documents
    
    @staticmethod
    def process_pptx(filepath: str) -> List[Dict]:
        """Extract text from PowerPoint"""
        prs = Presentation(filepath)
        documents = []
        
        for i, slide in enumerate(prs.slides):
            text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            
            text = '\n'.join(text_parts)
            if text.strip():
                documents.append({
                    'text': text,
                    'metadata': {
                        'source': os.path.basename(filepath),
                        'slide': i + 1,
                        'type': 'pptx'
                    }
                })
        
        return documents
    
    @staticmethod
    def process_file(filepath: str) -> List[Dict]:
        """Process any supported file format"""
        ext = os.path.splitext(filepath)[1].lower()
        
        processors = {
            '.pdf': DocumentProcessor.process_pdf,
            '.docx': DocumentProcessor.process_docx,
            '.txt': DocumentProcessor.process_txt,
            '.xlsx': DocumentProcessor.process_xlsx,
            '.pptx': DocumentProcessor.process_pptx,
        }
        
        if ext in processors:
            return processors[ext](filepath)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
